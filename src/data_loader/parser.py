"""
첨부파일 텍스트 추출
xlsx, xls, pdf, pptx, docx 에서 텍스트 추출
"""
import os
import logging
from pathlib import Path
from typing import List, Dict, Optional
import csv
import zipfile
import xml.etree.ElementTree as ET

PROJECT_ROOT = Path(__file__).parent.parent.parent
FILES_DIR = PROJECT_ROOT / "data" / "files"

logger = logging.getLogger(__name__)

# 지원 확장자
PARSEABLE = {".xlsx", ".xls", ".pdf", ".pptx", ".docx", ".txt", ".csv", ".hwpx", ".hwp"}
SKIP = {".gif", ".bmp", ".jpg", ".jpeg", ".png", ".ppt", ".doc"}


def parse_file(filepath: str) -> Optional[str]:
    """파일에서 텍스트 추출. 실패 시 None 반환."""
    path = Path(filepath)
    if not path.exists():
        return None

    ext = path.suffix.lower()
    if ext not in PARSEABLE:
        return None

    try:
        if ext == ".xlsx":
            return _parse_xlsx(path)
        if ext == ".xls":
            return _parse_xls(path)
        if ext == ".pdf":
            return _parse_pdf(path)
        if ext == ".pptx":
            return _parse_pptx(path)
        if ext == ".docx":
            return _parse_docx(path)
        if ext == ".txt":
            return _parse_txt(path)
        if ext == ".csv":
            return _parse_csv(path)
        if ext == ".hwpx":
            return _parse_hwpx(path)
        if ext == ".hwp":
            return _parse_hwp(path)
    except Exception:
        return None

    return None


def parse_post_files(post: Dict) -> str:
    """게시글의 모든 첨부파일 텍스트를 합쳐 반환"""
    texts = []
    post_id = str(post.get("id", ""))

    for f in post.get("files", []):
        local_path = f.get("local_path", "")
        if local_path and Path(local_path).exists():
            fpath = local_path
        else:
            # data/files/{post_id}/{filename} 에서 탐색
            name = f.get("name", "")
            candidates = []
            if name:
                candidates.append(FILES_DIR / post_id / name)
                # RAG 문서(id가 "sig_union:1856" 같은 형태) 지원:
                # data/files/<bo_table>/<wr_id>/<filename>
                if ":" in post_id:
                    prefix, rest = post_id.split(":", 1)
                    if prefix and rest:
                        candidates.append(FILES_DIR / prefix / rest / name)

            fpath = ""
            for c in candidates:
                if c and Path(c).exists():
                    fpath = str(c)
                    break
            if not fpath:
                continue

        ext = Path(fpath).suffix.lower()
        if ext in SKIP:
            continue

        content = parse_file(str(fpath))
        if content and content.strip():
            fname = f.get("name", Path(fpath).name)
            texts.append(f"[첨부: {fname}]\n{content}")

    return "\n\n".join(texts)


def enrich_posts_with_files(posts: List[Dict]) -> List[Dict]:
    """각 게시글에 file_content 필드를 추가"""
    total = len(posts)
    parsed = 0
    for i, p in enumerate(posts):
        if not p.get("files"):
            continue
        # 이미 파싱된 경우 재파싱 방지 (운영에서 RAG/게시판 데이터가 섞일 수 있음)
        if p.get("file_content"):
            continue
        file_text = parse_post_files(p)
        if file_text:
            p["file_content"] = file_text
            parsed += 1
        if (i + 1) % 50 == 0:
            logger.info("파일 파싱 진행: %d/%d", i + 1, total)

    logger.info("파일 파싱 완료: %d건 추출 (전체 %d건)", parsed, total)
    return posts


# ========== 개별 파서 ==========

def _parse_xlsx(path: Path) -> str:
    """xlsx 파일에서 모든 시트의 셀 값 추출"""
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        if len(wb.sheetnames) > 1:
            lines.append(f"[시트: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            row_text = "\t".join(cells).strip()
            if row_text:
                lines.append(row_text)
    wb.close()
    return "\n".join(lines)


def _parse_xls(path: Path) -> str:
    """xls 파일에서 텍스트 추출"""
    import xlrd
    wb = xlrd.open_workbook(str(path))
    lines = []
    for sheet in wb.sheets():
        if wb.nsheets > 1:
            lines.append(f"[시트: {sheet.name}]")
        for row_idx in range(sheet.nrows):
            cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
            row_text = "\t".join(cells).strip()
            if row_text:
                lines.append(row_text)
    return "\n".join(lines)


def _parse_pdf(path: Path) -> str:
    """PDF에서 텍스트 추출"""
    from PyPDF2 import PdfReader
    reader = PdfReader(str(path))
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text.strip())
    return "\n".join(texts)


def _parse_pptx(path: Path) -> str:
    """pptx에서 슬라이드 텍스트 추출"""
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            lines.append(f"[슬라이드 {i}]")
            lines.extend(slide_texts)
    return "\n".join(lines)


def _parse_docx(path: Path) -> str:
    """docx에서 텍스트 추출"""
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_txt(path: Path) -> str:
    """txt 텍스트 로드 (너무 큰 파일은 앞부분만)"""
    # 운영 데이터에 인코딩이 섞일 수 있어, 실패 시 cp949도 시도
    for enc in ("utf-8", "utf-8-sig", "cp949"):
        try:
            text = path.read_text(encoding=enc, errors="replace")
            return text[:200000]
        except Exception:
            continue
    return ""


def _parse_csv(path: Path) -> str:
    """csv 텍스트 추출 (표 -> 탭 구분 텍스트)"""
    raw = _parse_txt(path)
    if not raw:
        return ""
    reader = csv.reader(raw.splitlines())
    lines: List[str] = []
    for row in reader:
        line = "\t".join((c or "").strip() for c in row).strip()
        if line:
            lines.append(line)
        if sum(len(x) for x in lines) > 200000:
            break
    return "\n".join(lines)[:200000]


def _parse_hwpx(path: Path) -> str:
    """hwpx 텍스트 추출 (zip 내부 xml itertext)"""
    # hwpx는 zip 기반. 섹션 XML을 우선 추출.
    texts: List[str] = []
    with zipfile.ZipFile(str(path), "r") as zf:
        names = zf.namelist()
        section_xmls = [n for n in names if n.startswith("Contents/section") and n.endswith(".xml")]
        if not section_xmls:
            section_xmls = [n for n in names if n.endswith(".xml")]

        for name in section_xmls:
            try:
                data = zf.read(name)
                root = ET.fromstring(data)
                chunk = " ".join(t.strip() for t in root.itertext() if t and t.strip())
                if chunk:
                    texts.append(chunk)
                if sum(len(x) for x in texts) > 200000:
                    break
            except Exception:
                continue
    return "\n".join(texts)[:200000]


def _parse_hwp(path: Path) -> str:
    """hwp 텍스트 추출 (OLE PrvText 미리보기 스트림 기반)"""
    # HWP v5는 OLE compound file이며, PrvText 스트림이 UTF-16LE로 들어있는 경우가 많음.
    # (전체 본문 변환은 별도 라이브러리/툴 필요하지만, 미리보기 텍스트만으로도 검색 품질이 개선됨)
    try:
        import olefile  # type: ignore
    except Exception:
        return ""

    try:
        with olefile.OleFileIO(str(path)) as ole:
            # 후보 스트림명들
            candidates = [
                ["PrvText.utf8"],
                ["PrvText"],
                ["PrvText", "utf8"],
            ]
            data = b""
            for c in candidates:
                try:
                    if ole.exists(c):
                        data = ole.openstream(c).read()
                        if data:
                            break
                except Exception:
                    continue

        if not data:
            return ""

        # utf-16le가 대부분. 실패 시 utf-8도 시도.
        for enc in ("utf-16le", "utf-8", "cp949"):
            try:
                text = data.decode(enc, errors="replace")
                return text.strip()[:200000]
            except Exception:
                continue
    except Exception:
        return ""

    return ""
